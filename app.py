# -*- coding: utf-8 -*-
import streamlit as st
import requests
import os
import zipfile
import io
import xml.etree.ElementTree as ET
from google import genai
from google.genai import types
from pydantic import BaseModel
import datetime
import json
import re
import pypdf
import docx
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE

# Safety Default
target_standard = locals().get('target_standard') or globals().get('target_standard') or "General Standard"

# --- Configuration ---
st.set_page_config(page_title="Canvas Content Creator", page_icon="✨", layout="wide", initial_sidebar_state="expanded")

# Initialize Session State
if 'generated_prompt' not in st.session_state:
    st.session_state['generated_prompt'] = ""
if 'lesson_plan_pdf' not in st.session_state:
    st.session_state['lesson_plan_pdf'] = None
if 'slide_deck_pptx' not in st.session_state:
    st.session_state['slide_deck_pptx'] = None
if 'is_generated' not in st.session_state:
    st.session_state['is_generated'] = False
if 'selected_tool_name' not in st.session_state:
    st.session_state['selected_tool_name'] = "None"

# Custom CSS matching AI Teacher Lounge branding
st.markdown("""
<style>
/* ===================================
   AI Teacher Lounge Brand Theme
   Colors: Dark (#3C453C) / White (#FFFFFF)
   Accent: Green (#5a9a5a)
   =================================== */

/* FORCE SIDEBAR TO STAY OPEN - Multiple approaches */
[data-testid="collapsedControl"] {
    display: none !important;
    visibility: hidden !important;
    width: 0 !important;
    height: 0 !important;
}

[data-testid="stSidebar"][aria-expanded="false"] {
    display: block !important;
    width: 21rem !important;
    min-width: 21rem !important;
    transform: none !important;
    margin-left: 0 !important;
}

[data-testid="stSidebar"] {
    width: 21rem !important;
    min-width: 21rem !important;
}

/* Hide collapse button arrow */
button[kind="header"] {
    display: none !important;
}

/* Hide any "key" or icon fallback text */
[data-testid="stSidebar"] span.material-icons,
[data-testid="stSidebar"] span[class*="icon"] {
    font-size: 0 !important;
    visibility: hidden !important;
}

/* Hide specific "key" text that appears */
[data-testid="stSidebar"] > div > div > div > div:first-child {
    overflow: hidden !important;
}

/* Import Google Fonts */
@import url('https://fonts.googleapis.com/css2?family=Open+Sans:wght@400;500;600;700&family=Merriweather:wght@400;700&display=swap');

/* CSS Variables */
:root {
    --brand-dark: #3C453C;
    --brand-dark-hover: #2d332d;
    --brand-dark-light: #4a554a;
    --accent: #5a9a5a;
    --accent-light: #7ab87a;
    --accent-dark: #3d6b3d;
    --bg-light: #FFFFFF;
    --bg-off-white: #f8faf8;
    --bg-cream: #f5f7f5;
    --text-dark: #3C453C;
    --text-medium: #5a635a;
    --text-light: #7a837a;
    --border-light: #e5e8e5;
    --border-medium: #d0d5d0;
    --shadow-sm: 0 2px 8px rgba(60, 69, 60, 0.08);
    --shadow-md: 0 4px 16px rgba(60, 69, 60, 0.12);
    --radius-sm: 6px;
    --radius-md: 12px;
    --radius-lg: 20px;
    --radius-full: 100px;
}

/* Main App Background */
.stApp {
    background-color: var(--bg-off-white) !important;
}

/* Main content area */
.main .block-container {
    background-color: var(--bg-light);
    border-radius: var(--radius-lg);
    padding: 2rem 2rem 3rem 2rem !important;
    margin: 1rem;
    box-shadow: var(--shadow-md);
    max-width: 1200px;
}

/* Typography */
h1, h2, h3, h4, h5, h6 {
    font-family: 'Merriweather', serif !important;
    color: var(--text-dark) !important;
    font-weight: 700 !important;
}

p, span, label, div {
    font-family: 'Open Sans', sans-serif !important;
}

/* Main Title Styling */
.main h1 {
    color: var(--brand-dark) !important;
    padding-bottom: 0.75rem;
    border-bottom: 3px solid var(--accent);
    margin-bottom: 0.5rem !important;
    font-size: 2.25rem !important;
}

/* Caption under title */
.main [data-testid="stCaptionContainer"] {
    color: var(--text-medium) !important;
    font-size: 1rem !important;
    margin-bottom: 1rem !important;
}

/* Section Headers */
h2 {
    color: var(--brand-dark) !important;
    font-size: 1.5rem !important;
    margin-top: 1.5rem !important;
    margin-bottom: 1rem !important;
}

h3 {
    color: var(--text-dark) !important;
    font-size: 1.125rem !important;
}

/* Sidebar Styling */
[data-testid="stSidebar"] {
    background-color: #3C453C !important;
    border-right: none !important;
}

/* Hide the "key" text and any icon fallbacks at top of sidebar */
[data-testid="stSidebar"] > div:first-child > div:first-child > div:first-child {
    font-size: 0 !important;
    height: 0 !important;
    overflow: hidden !important;
    visibility: hidden !important;
}

/* Hide any standalone text that says "key" or similar icon fallback */
[data-testid="stSidebarContent"] > div:first-child:not([data-testid]) {
    display: none !important;
}

[data-testid="stSidebar"] * {
    color: white !important;
    font-family: 'Open Sans', sans-serif !important;
}

[data-testid="stSidebar"] h1,
[data-testid="stSidebar"] h2,
[data-testid="stSidebar"] h3 {
    color: white !important;
    border-bottom: 2px solid #5a9a5a !important;
    padding-bottom: 0.5rem;
    font-family: 'Merriweather', serif !important;
}

[data-testid="stSidebar"] label {
    color: white !important;
    font-weight: 500 !important;
}

/* File uploader label - specific color #101217 */
[data-testid="stSidebar"] [data-testid="stFileUploader"] label {
    color: #101217 !important;
}

[data-testid="stSidebar"] [data-testid="stFileUploader"] label p {
    color: #101217 !important;
}

[data-testid="stSidebar"] [data-testid="stFileUploader"] [data-testid="stWidgetLabel"] p {
    color: #101217 !important;
}

[data-testid="stSidebar"] [data-testid="stFileUploader"] section {
    background-color: rgba(255, 255, 255, 0.9) !important;
    border-radius: 6px !important;
}

[data-testid="stSidebar"] [data-testid="stFileUploader"] section small {
    color: #101217 !important;
}

[data-testid="stSidebar"] [data-testid="stFileUploader"] section span {
    color: #101217 !important;
}

/* Sidebar Input Fields */
[data-testid="stSidebar"] input,
[data-testid="stSidebar"] textarea,
[data-testid="stSidebar"] [data-baseweb="select"] > div {
    background-color: rgba(255, 255, 255, 0.1) !important;
    border: 1px solid rgba(255, 255, 255, 0.2) !important;
    border-radius: 6px !important;
    color: white !important;
    font-family: 'Open Sans', sans-serif !important;
}

[data-testid="stSidebar"] input:focus,
[data-testid="stSidebar"] textarea:focus {
    border-color: #5a9a5a !important;
}

/* Sidebar subheaders */
[data-testid="stSidebar"] [data-testid="stSubheader"] {
    color: white !important;
    font-size: 0.95rem !important;
    margin-top: 0.5rem !important;
}

/* Sidebar dividers */
[data-testid="stSidebar"] hr {
    border-color: rgba(255, 255, 255, 0.2) !important;
    margin: 1rem 0 !important;
}

/* Main Content Input Fields */
.main input,
.main textarea,
.main [data-baseweb="select"] > div {
    background-color: var(--bg-light) !important;
    border: 1px solid var(--border-medium) !important;
    border-radius: var(--radius-sm) !important;
    color: var(--text-dark) !important;
    font-family: 'Open Sans', sans-serif !important;
}

.main input:focus,
.main textarea:focus {
    border-color: var(--brand-dark) !important;
    box-shadow: 0 0 0 2px rgba(60, 69, 60, 0.15) !important;
}

/* Text Area Styling */
.stTextArea textarea {
    background-color: var(--bg-cream) !important;
    color: var(--text-dark) !important;
    font-family: 'Courier New', monospace !important;
    border: 1px solid var(--border-medium) !important;
    border-radius: var(--radius-md) !important;
    padding: 1rem !important;
}

/* Primary Buttons */
.stButton > button {
    background-color: var(--brand-dark) !important;
    color: white !important;
    border: 2px solid var(--brand-dark) !important;
    border-radius: var(--radius-full) !important;
    padding: 0.6rem 1.5rem !important;
    font-weight: 600 !important;
    font-family: 'Open Sans', sans-serif !important;
    transition: all 0.3s ease !important;
}

.stButton > button:hover {
    background-color: var(--brand-dark-hover) !important;
    border-color: var(--brand-dark-hover) !important;
    transform: translateY(-2px) !important;
    box-shadow: var(--shadow-md) !important;
}

/* Download Buttons */
.stDownloadButton > button {
    background-color: var(--accent) !important;
    color: white !important;
    border: 2px solid var(--accent) !important;
    border-radius: var(--radius-full) !important;
    padding: 0.6rem 1.5rem !important;
    font-weight: 600 !important;
    transition: all 0.3s ease !important;
}

.stDownloadButton > button:hover {
    background-color: var(--accent-dark) !important;
    border-color: var(--accent-dark) !important;
    transform: translateY(-2px) !important;
}

/* Checkbox and Toggle Styling */
[data-testid="stSidebar"] .stCheckbox label,
[data-testid="stSidebar"] .stCheckbox span,
[data-testid="stSidebar"] .stCheckbox p {
    color: white !important;
}

/* Main Content Checkbox - FORCE labels to show */
.main .stCheckbox {
    color: #3C453C !important;
}

.main .stCheckbox label {
    color: #3C453C !important;
    display: flex !important;
    visibility: visible !important;
}

.main .stCheckbox label > span {
    color: #3C453C !important;
    visibility: visible !important;
    display: inline !important;
}

.main .stCheckbox label p {
    color: #3C453C !important;
    visibility: visible !important;
    display: inline !important;
    opacity: 1 !important;
}

.main .stCheckbox [data-testid="stWidgetLabel"] {
    color: #3C453C !important;
    visibility: visible !important;
    display: block !important;
}

.main .stCheckbox [data-testid="stWidgetLabel"] p {
    color: #3C453C !important;
    visibility: visible !important;
    opacity: 1 !important;
}

.main .stCheckbox [data-testid="stMarkdownContainer"] {
    color: #3C453C !important;
    visibility: visible !important;
}

.main .stCheckbox [data-testid="stMarkdownContainer"] p {
    color: #3C453C !important;
    visibility: visible !important;
    opacity: 1 !important;
}

/* Toggle Switch */
[data-testid="stToggle"] > label > div[data-checked="true"] {
    background-color: #5a9a5a !important;
}

/* Expander in Main Content */
.main .streamlit-expanderHeader {
    background-color: var(--bg-cream) !important;
    border-radius: var(--radius-sm) !important;
    color: var(--text-dark) !important;
    font-weight: 600 !important;
    border: 1px solid var(--border-light) !important;
}

.main .streamlit-expanderHeader:hover {
    background-color: var(--bg-off-white) !important;
    border-color: var(--brand-dark) !important;
}

.main .streamlit-expanderContent {
    border: 1px solid var(--border-light) !important;
    border-top: none !important;
    border-radius: 0 0 var(--radius-sm) var(--radius-sm) !important;
    background-color: var(--bg-light) !important;
}

/* Info Box */
.stAlert {
    background-color: var(--bg-cream) !important;
    border: 1px solid var(--border-light) !important;
    border-radius: var(--radius-md) !important;
    color: var(--text-dark) !important;
}

.stAlert > div {
    color: var(--text-medium) !important;
}

/* Success Message */
[data-testid="stAlert"][data-baseweb="notification"] {
    border-radius: var(--radius-md) !important;
}

/* Progress Bar */
.stProgress > div > div {
    background-color: var(--accent) !important;
}

/* Selectbox Dropdown */
[data-baseweb="popover"] {
    border-radius: var(--radius-sm) !important;
}

[data-baseweb="menu"] {
    background-color: var(--bg-light) !important;
    border: 1px solid var(--border-light) !important;
}

[data-baseweb="menu"] li {
    color: var(--text-dark) !important;
}

[data-baseweb="menu"] li:hover {
    background-color: var(--bg-cream) !important;
}

/* Multiselect Tags */
[data-baseweb="tag"] {
    background-color: var(--brand-dark) !important;
    border-radius: var(--radius-full) !important;
}

/* Number Input */
.stNumberInput input {
    border-radius: var(--radius-sm) !important;
}

/* Date/Time Input */
.stDateInput input,
.stTimeInput input {
    border-radius: var(--radius-sm) !important;
}

/* Divider */
hr {
    border-color: var(--border-light) !important;
    margin: 1.5rem 0 !important;
}

/* Code Block */
code {
    background-color: var(--bg-cream) !important;
    color: var(--accent-dark) !important;
    padding: 0.2rem 0.4rem !important;
    border-radius: var(--radius-sm) !important;
    font-family: 'Courier New', monospace !important;
}

/* Markdown Text */
.stMarkdown p {
    color: var(--text-medium) !important;
    line-height: 1.6 !important;
}

/* Spinner */
.stSpinner > div {
    border-color: var(--accent) transparent transparent transparent !important;
}

/* File Uploader */
[data-testid="stFileUploader"] {
    border: 2px dashed var(--border-medium) !important;
    border-radius: var(--radius-md) !important;
    background-color: var(--bg-cream) !important;
}

[data-testid="stFileUploader"]:hover {
    border-color: var(--brand-dark) !important;
}

/* Success/Error States */
.element-container .stSuccess {
    background-color: rgba(90, 154, 90, 0.1) !important;
    border: 1px solid var(--accent) !important;
    color: var(--accent-dark) !important;
}

.element-container .stError {
    background-color: rgba(248, 113, 113, 0.1) !important;
    border: 1px solid #f87171 !important;
}

/* Hide Streamlit Branding */
#MainMenu {visibility: hidden;}
footer {visibility: hidden;}
header {visibility: hidden;}

/* Scrollbar Styling */
::-webkit-scrollbar {
    width: 8px;
    height: 8px;
}

::-webkit-scrollbar-track {
    background: var(--bg-cream);
    border-radius: 4px;
}

::-webkit-scrollbar-thumb {
    background: var(--brand-dark-light);
    border-radius: 4px;
}

::-webkit-scrollbar-thumb:hover {
    background: var(--brand-dark);
}

/* Responsive Adjustments */
@media (max-width: 768px) {
    .main .block-container {
        padding: 1rem !important;
        margin: 0.5rem;
    }
}
</style>
""", unsafe_allow_html=True)

# --- Constants ---

STEM_TOOLS = {
    "PhET: Balancing Chemical Equations": "https://phet.colorado.edu/sims/html/balancing-chemical-equations/latest/balancing-chemical-equations_en.html",
    "PhET: Circuit Construction Kit": "https://phet.colorado.edu/sims/html/circuit-construction-kit-dc/latest/circuit-construction-kit-dc_en.html",
    "PhET: Energy Skate Park": "https://phet.colorado.edu/sims/html/energy-skate-park/latest/energy-skate-park_en.html",
    "PhET: Natural Selection": "https://phet.colorado.edu/sims/html/natural-selection/latest/natural-selection_en.html",
    "PhET: Projectile Motion": "https://phet.colorado.edu/sims/html/projectile-motion/latest/projectile-motion_en.html",
    "PhET: Forces and Motion": "https://phet.colorado.edu/sims/html/forces-and-motion-basics/latest/forces-and-motion-basics_en.html",
    "Desmos: Graphing Calculator": "https://www.desmos.com/calculator",
    "Desmos: Scientific Calculator": "https://www.desmos.com/scientific",
    "GeoGebra: Geometry": "https://www.geogebra.org/geometry",
    "YouTube: Crash Course": "https://www.youtube.com/user/crashcourse",
    "YouTube: Khan Academy": "https://www.youtube.com/user/khanacademy",
    "YouTube: National Geographic": "https://www.youtube.com/user/NationalGeographic",
    "Wikipedia": "https://www.wikipedia.org/",
    "Google Slides": "https://docs.google.com/presentation/u/0/",
    "Canva": "https://www.canva.com/",
    "Desmos: Supply & Demand Shifters": "https://www.desmos.com/calculator/6mmm8psho7",
    "EconGraphs: Competitive Market": "https://www.econgraphs.org/graphs/micro/equilibrium/supply_and_demand_old",
    "Marginal Revolution: Elasticity Practice": "https://practice.mru.org/interactive-practice-supply-and-demand/",
    "Omni Margin Calculator": "https://www.omnicalculator.com/finance/margin",
    "AutoDraw": "https://www.autodraw.com/",
    "Sketchpad": "https://sketch.io/sketchpad/",
    "Color Wheel": "https://color.adobe.com/create/color-wheel",
    "Google Arts & Culture": "https://artsandculture.google.com/",
    "Python Online Compiler": "https://trinket.io/embed/python3",
    "Scratch": "https://scratch.mit.edu/projects/editor/embed"
}

# --- Helper Functions (QTI) ---

class Question(BaseModel):
    question_text: str
    options: list[str]
    correct_answer_index: int

class Quiz(BaseModel):
    questions: list[Question]

def create_imsmanifest():
    """Creates the imsmanifest.xml content."""
    manifest_template = """<?xml version="1.0" encoding="UTF-8"?>
<manifest identifier="man00001" xmlns="http://www.imsglobal.org/xsd/imscp_v1p1" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance" xsi:schemaLocation="http://www.imsglobal.org/xsd/imscp_v1p1 http://www.imsglobal.org/xsd/imscp_v1p1.xsd">
  <metadata>
    <schema>IMS Content</schema>
    <schemaversion>1.1.3</schemaversion>
  </metadata>
  <organizations/>
  <resources>
    <resource identifier="res00001" type="imsqti_xmlv1p2">
      <file href="quiz.xml"/>
    </resource>
  </resources>
</manifest>"""
    return manifest_template.encode('utf-8')

def create_quiz_xml(questions, title="Generated Quiz"):
    """Creates the quiz.xml content (QTI v1.2)."""
    root = ET.Element("questestinterop", {
        "xmlns": "http://www.imsglobal.org/xsd/ims_qtiasiv1p2",
        "xmlns:xsi": "http://www.w3.org/2001/XMLSchema-instance",
        "xsi:schemaLocation": "http://www.imsglobal.org/xsd/ims_qtiasiv1p2 http://www.imsglobal.org/xsd/ims_qtiasiv1p2.xsd"
    })

    assessment = ET.SubElement(root, "assessment", {
        "ident": "quiz001",
        "title": title
    })

    section = ET.SubElement(assessment, "section", {
        "ident": "sec001",
        "title": "Main Section"
    })

    for i, q in enumerate(questions):
        item = ET.SubElement(section, "item", {
            "ident": f"q{i+1}",
            "title": f"Question {i+1}"
        })

        # Question Text
        presentation = ET.SubElement(item, "presentation")
        material = ET.SubElement(presentation, "material")
        mattext = ET.SubElement(material, "mattext", {"texttype": "text/html"})
        # Use placeholder for CDATA
        mattext.text = f"__CDATA_START__{q.question_text}__CDATA_END__"

        # Response Lid (Multiple Choice)
        response_lid = ET.SubElement(presentation, "response_lid", {
            "ident": f"response_{i+1}",
            "rcardinality": "Single"
        })
        render_choice = ET.SubElement(response_lid, "render_choice")

        for j, option in enumerate(q.options):
            response_label = ET.SubElement(render_choice, "response_label", {"ident": f"opt_{i+1}_{j}"})
            material_opt = ET.SubElement(response_label, "material")
            mattext_opt = ET.SubElement(material_opt, "mattext", {"texttype": "text/html"})
            mattext_opt.text = f"__CDATA_START__{option}__CDATA_END__"

        # Processing (Correct Answer)
        resprocessing = ET.SubElement(item, "resprocessing")
        outcomes = ET.SubElement(resprocessing, "outcomes")
        decvar = ET.SubElement(outcomes, "decvar", {
            "defaultval": "0",
            "varname": "SCORE",
            "vartype": "Integer"
        })

        respcondition = ET.SubElement(resprocessing, "respcondition", {"continue": "No"})
        conditionvar = ET.SubElement(respcondition, "conditionvar")
        correct_ident = f"opt_{i+1}_{q.correct_answer_index}"
        varequal = ET.SubElement(conditionvar, "varequal", {"respident": f"response_{i+1}"})
        varequal.text = correct_ident

        setvar = ET.SubElement(respcondition, "setvar", {
            "action": "Set",
            "varname": "SCORE"
        })
        setvar.text = "1"

    return ET.tostring(root, encoding='utf-8', xml_declaration=True)

def generate_qti_zip(quiz_data, title="Generated Quiz"):
    """Generates a QTI 1.2 Zip package."""
    try:
        # 1. Create Manifest
        manifest_data = create_imsmanifest()
        
        # 2. Create Quiz XML
        questions = [Question(**q) for q in quiz_data['questions']]
        quiz_xml_raw = create_quiz_xml(questions, title=title)
        
        # 3. Post-Process for CDATA
        # Decode bytes to string for replacement
        xml_str = quiz_xml_raw.decode('utf-8')
        # Replace placeholders with actual CDATA tags
        # Note: We also need to unescape the content inside, as ElementTree might have escaped < > &
        # But since we used placeholders, the content inside is just text. 
        # The key is that we want the FINAL output to have <![CDATA[ ... ]]>
        
        xml_str = xml_str.replace("__CDATA_START__", "<![CDATA[").replace("__CDATA_END__", "]]>")
        
        # Re-encode
        final_quiz_xml = xml_str.encode('utf-8')
        
        # 4. Zip It
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("imsmanifest.xml", manifest_data)
            zf.writestr("quiz.xml", final_quiz_xml)
            
        return zip_buffer.getvalue()
    except Exception as e:
        st.error(f"Error creating QTI Zip: {e}")
        return None

def generate_quiz_json(prompt):
    """Generates the Quiz JSON data from Gemini."""
    client = get_gemini_client()
    if not client: return None
    
    try:
        response = client.models.generate_content(
            model='gemini-2.0-flash',
            contents=prompt,
            config=types.GenerateContentConfig(response_mime_type='application/json')
        )
        return clean_json(response.text)
    except Exception as e:
        st.error(f"Error generating quiz JSON: {e}")
        return None

def generate_quiz_data_batched(topic, subtopic, target_count, due_date, due_time, points_per_question, question_types, grade_level, is_sped, is_gifted, is_ml, language, source_text="", context_topics=None):
    """Generates quiz questions in batches to ensure target count is met."""
    all_questions = []
    batch_size = 10
    # Over-provision: (target // 10) + 2 batches
    num_batches = (target_count // 10) + 2
    
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    for i in range(num_batches):
        status_text.text(f"Generating Batch {i+1}/{num_batches}...")
        
        # Construct prompt for this batch
        # Note: We use batch_size here, not target_count
        prompt = construct_quiz_prompt(topic, subtopic, batch_size, due_date, due_time, points_per_question, question_types, grade_level, is_sped, is_gifted, is_ml, language, source_text, context_topics)
        
        # Generate JSON
        batch_data = generate_quiz_json(prompt)
        
        if batch_data and 'questions' in batch_data:
            all_questions.extend(batch_data['questions'])
        
        # Update progress
        progress_bar.progress((i + 1) / num_batches)
        
        # Early break if we have enough
        if len(all_questions) >= target_count:
            break
            
    status_text.empty()
    progress_bar.empty()
    
    # Trim to exact count
    all_questions = all_questions[:target_count]
    
    return {"questions": all_questions}

def generate_unit_sequence_json(prompt):
    """Generates the Unit Sequence JSON from Gemini."""
    client = get_gemini_client()
    if not client: return None
    
    try:
        response = client.models.generate_content(
            model='gemini-2.0-flash',
            contents=prompt,
            config=types.GenerateContentConfig(response_mime_type='application/json')
        )
        return clean_json(response.text)
    except Exception as e:
        st.error(f"Error generating unit sequence: {e}")
        return None

def generate_unit_package(sequence_data, topic, grade_level, is_sped, is_gifted, is_ml, language, subject, strategy, source_text, due_date, due_time, points, points_per_question, question_types, standard="General Standard"):
    """Generates all files for a unit and zips them, including Lesson Plans and Slides for each Assignment."""
    context_buffer = []
    zip_buffer = io.BytesIO()
    
    progress_bar = st.progress(0)
    status_text = st.empty()
    total_items = len(sequence_data)
    
    # Calculate progress steps: each Assignment has 3 sub-steps (HTML, PDF, PPTX), each Quiz has 1
    total_steps = sum(3 if item.get('type') == 'Assignment' else 1 for item in sequence_data)
    current_step = 0
    
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, item in enumerate(sequence_data):
            idx = i + 1
            item_type = item.get('type')
            title = item.get('title', f"Item {idx}")
            focus = item.get('focus_topic', topic)
            safe_title = title.replace(" ", "_").replace("/", "-")
            
            if item_type == "Assignment":
                # --- Step 1: Generate Assignment HTML ---
                status_text.text(f"Generating Assignment {idx}/{total_items}: {title} (HTML)...")
                prompt = construct_assignment_prompt(topic, focus, "None", due_date, due_time, points, grade_level, is_sped, is_gifted, is_ml, language, subject, strategy, source_text)
                
                client = get_gemini_client()
                if client:
                    try:
                        response = client.models.generate_content(
                            model='gemini-2.0-flash',
                            contents=prompt
                        )
                        html_content = response.text
                        # Clean markdown code blocks if present
                        if html_content.startswith("```html"):
                            html_content = html_content[7:]
                        if html_content.endswith("```"):
                            html_content = html_content[:-3]
                            
                        zf.writestr(f"{idx:02d}_Assignment_{safe_title}.html", html_content)
                        context_buffer.append(focus)
                    except Exception as e:
                        print(f"Error generating assignment {title}: {e}")
                
                current_step += 1
                progress_bar.progress(current_step / total_steps)
                
                # --- Step 2: Generate Lesson Plan PDF (with focus-specific content) ---
                status_text.text(f"Generating Assignment {idx}/{total_items}: {title} (Lesson Plan)...")
                lesson_plan_pdf, lesson_plan_text = generate_lesson_plan_pdf(
                    topic=f"{topic}: {focus}",  # Include subtopic for specificity
                    standard=standard,
                    grade=grade_level,
                    strategy=strategy
                )
                
                if lesson_plan_pdf:
                    zf.writestr(f"{idx:02d}_LessonPlan_{safe_title}.pdf", lesson_plan_pdf)
                
                current_step += 1
                progress_bar.progress(current_step / total_steps)
                
                # --- Step 3: Generate Slide Deck (CHAINED from Lesson Plan to prevent overlap) ---
                status_text.text(f"Generating Assignment {idx}/{total_items}: {title} (Slides)...")
                # CRITICAL: Pass lesson_plan_text as source_text to ensure slides are
                # complementary (keywords only) and don't duplicate lesson plan content
                slide_deck_pptx = generate_slide_deck(
                    topic=f"{topic}: {focus}",
                    grade=grade_level,
                    strategy=strategy,
                    source_text=lesson_plan_text if lesson_plan_text else ""
                )
                
                if slide_deck_pptx:
                    zf.writestr(f"{idx:02d}_Slides_{safe_title}.pptx", slide_deck_pptx)
                
                current_step += 1
                progress_bar.progress(current_step / total_steps)
                        
            elif item_type == "Quiz":
                # Generate Quiz Zip
                status_text.text(f"Generating Quiz {idx}/{total_items}: {title}...")
                # Use context_buffer for contextual awareness
                # Default to 10 questions for unit quizzes
                quiz_data = generate_quiz_data_batched(topic, focus, 10, due_date, due_time, points_per_question, question_types, grade_level, is_sped, is_gifted, is_ml, language, source_text, context_topics=context_buffer)
                
                if quiz_data:
                    qti_zip = generate_qti_zip(quiz_data, title=title)
                    if qti_zip:
                        zf.writestr(f"{idx:02d}_Quiz_{safe_title}.zip", qti_zip)
                
                # Clear context after quiz
                context_buffer = []
                
                current_step += 1
                progress_bar.progress(current_step / total_steps)
            
    status_text.empty()
    progress_bar.empty()
    return zip_buffer.getvalue()

# --- AI Generation Functions ---

def get_gemini_client():
    try:
        api_key = st.secrets.get("GEMINI_API_KEY")
        if not api_key:
            st.error("⚠️ GEMINI_API_KEY not configured in secrets")
            st.stop()
        return genai.Client(api_key=api_key)
    except Exception as e:
        st.error(f"Failed to initialize Gemini client: {e}")
        st.stop()
        return None

def check_api_connection():
    """Checks connection to Gemini API."""
    try:
        api_key = st.secrets["GEMINI_API_KEY"]
        url = f"https://generativelanguage.googleapis.com/v1beta/models?key={api_key}"
        response = requests.get(url)
        
        if response.status_code != 200:
            st.error(f"Connection Error: {response.status_code}")
        else:
            st.toast("Connected successfully!", icon="✅")
    except Exception as e:
        st.error(f"Connection Failed: {e}")
        st.error("GEMINI_API_KEY not found in secrets.toml")
        return None

def clean_json(text):
    """Cleans JSON text by removing markdown fences and fixing unescaped backslashes."""
    text = text.strip()
    if text.startswith("```json"):
        text = text[7:]
    if text.endswith("```"):
        text = text[:-3]
    
    # Fix unescaped backslashes (common in LaTeX)
    # Find backslashes that are NOT followed by valid escape chars (n, r, t, b, f, u, ", \, /)
    # and double them up.
    text = re.sub(r'\\(?![nrtbfu"/\\\\])', r'\\\\', text)
    
    return json.loads(text)

def recommend_tool(topic, standard):
    client = get_gemini_client()
    if not client: return "None"

    tools_keys = list(STEM_TOOLS.keys())
    prompt = f"""
    Given the topic "{topic}" and standard "{standard}", which ONE of these tools is the absolute best match?
    Options: {tools_keys}
    
    Return ONLY the exact dictionary key. If nothing fits perfectly, return "None".
    """
    try:
        response = client.models.generate_content(
            model='gemini-2.0-flash',
            contents=prompt
        )
        recommended = response.text.strip()
        if recommended in tools_keys:
            return recommended
        return "None"
    except Exception as e:
        st.error(f"Error recommending tool: {e}")
        return "None"

def generate_unit_outline(topic, num_assignments, num_quizzes):
    # Safety Default
    target_standard = locals().get('target_standard') or globals().get('target_standard') or "General Standard"
    
    client = get_gemini_client()
    if not client: return []

    prompt = f"""
    Create a unit outline for the topic: {topic}.
    Generate exactly {num_assignments} assignment titles and {num_quizzes} quiz titles.
    Return a JSON object with a list 'items', where each item has 'type' ('Assignment' or 'Quiz') and 'title'.
    """
    try:
        response = client.models.generate_content(
            model='gemini-2.0-flash',
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type='application/json'
            )
        )
        return clean_json(response.text)['items']
    except Exception as e:
        st.error(f"Error generating outline: {e}")
        return []

def generate_lesson_plan_pdf(topic, standard, grade, strategy="None / Standard"):
    """Generates a High-Design 5E Lesson Plan PDF (Strict One-Page). Returns (pdf_bytes, raw_text)."""
    # Safety Default
    target_standard = locals().get('target_standard') or globals().get('target_standard') or "General Standard"

    client = get_gemini_client()
    if not client: return None, ""

    # 1. AI Generation (Structured JSON)
    prompt = f"""
    Create a 5E Lesson Plan for Grade {grade} on "{topic}" (Standard: {standard}).
    Return a JSON object with this EXACT structure:
    {{
        "metadata": {{
            "duration": "e.g., 60 minutes",
            "materials": ["item 1", "item 2"],
            "vocabulary": ["term 1", "term 2"],
            "differentiation": {{
                "sped": ["mod 1", "mod 2"],
                "ml": ["support 1", "support 2"]
            }}
        }},
        "sections": [
            {{"phase": "Engage", "time": "10 mins", "activity": "Brief description..."}},
            {{"phase": "Explore", "time": "15 mins", "activity": "Brief description..."}},
            {{"phase": "Explain", "time": "10 mins", "activity": "Brief description..."}},
            {{"phase": "Elaborate", "time": "15 mins", "activity": "Independent Practice referencing the {topic} Assignment..."}},
            {{"phase": "Evaluate", "time": "10 mins", "activity": "Assessment referencing the {topic} Quiz..."}}
        ]
    }}
    Keep descriptions concise (bullet points preferred).
    """
    
    try:
        response = client.models.generate_content(
            model='gemini-2.0-flash',
            contents=prompt,
            config=types.GenerateContentConfig(response_mime_type='application/json')
        )
        data = clean_json(response.text)
        
        # 2. PDF Creation (High-Design Dashboard - Strict One Page)
        pdf = FPDF(orientation='P', unit='mm', format='Letter')
        pdf.add_page()
        pdf.set_auto_page_break(auto=False) # Disable auto page break
        
        # Colors
        header_bg = (30, 36, 58) # Dark Blue
        sidebar_bg = (240, 244, 248) # Light Grey/Blue
        text_color = (0, 0, 0)
        header_text_color = (255, 255, 255)
        
        # --- Header (Full Width) ---
        pdf.set_fill_color(*header_bg)
        pdf.rect(0, 0, 216, 38, 'F') # 1.5 inch approx 38mm
        
        pdf.set_text_color(*header_text_color)
        pdf.set_font("Helvetica", 'B', 16)
        pdf.set_xy(10, 10)
        pdf.cell(0, 8, f"Lesson Plan: {topic}", ln=1)
        
        # Strategy (Top Right)
        pdf.set_xy(120, 10)
        pdf.set_font("Helvetica", 'I', 10)
        pdf.cell(86, 8, f"Strategy: {strategy}", ln=1, align='R')
        
        pdf.set_font("Helvetica", '', 10)
        pdf.set_text_color(200, 200, 200) # Light Grey
        pdf.set_xy(10, 20)
        pdf.cell(0, 5, f"Grade: {grade}", ln=1)
        
        # Truncate Standard
        std_desc = standard
        if len(std_desc) > 120:
            std_desc = std_desc[:117] + "..."
        pdf.multi_cell(0, 5, f"Standard: {std_desc}")
        
        # --- Sidebar (Left 2.5 inches -> 63.5mm) ---
        sidebar_width = 64
        pdf.set_fill_color(*sidebar_bg)
        pdf.rect(0, 38, sidebar_width, 241, 'F')
        
        # Sidebar Content
        pdf.set_text_color(*text_color)
        y_pos = 45
        x_pos = 5
        
        def sidebar_section(title, items):
            nonlocal y_pos
            if y_pos > 250: return # Stop if too low
            
            # Step A: Set XY
            pdf.set_xy(x_pos, y_pos)
            
            # Step B: Print Title
            pdf.set_font("Helvetica", 'B', 9)
            pdf.cell(50, 5, title.upper(), ln=1)
            
            # Step C: Print Content
            pdf.set_font("Helvetica", '', 8)
            content_str = ""
            if isinstance(items, list):
                for item in items:
                    content_str += f"- {item}\n"
            elif isinstance(items, str):
                content_str = items
            
            safe_content = content_str.encode('latin-1', 'replace').decode('latin-1')
            
            # Save current Y before printing content? No, we print then check Y.
            # Actually, we need to set XY for content? No, ln=1 moved us down.
            # But let's be precise as requested: "Step C: Save the current Y. Print the Content"
            # The multi_cell handles the printing.
            pdf.set_xy(x_pos, pdf.get_y()) 
            pdf.multi_cell(55, 4, safe_content)
            
            # Step D: Update current_y
            y_pos = pdf.get_y() + 10

        sidebar_section("Duration", data['metadata'].get('duration', '60 mins'))
        sidebar_section("Materials", data['metadata'].get('materials', []))
        sidebar_section("Vocabulary", data['metadata'].get('vocabulary', []))
        
        # Differentiation
        if y_pos < 250:
            pdf.set_xy(x_pos, y_pos)
            pdf.set_font("Helvetica", 'B', 9)
            pdf.cell(50, 5, "DIFFERENTIATION", ln=1)
            # Update Y for content
            y_pos = pdf.get_y()
            
            diff = data['metadata'].get('differentiation', {})
            if diff.get('sped'):
                pdf.set_xy(x_pos, y_pos)
                pdf.set_font("Helvetica", 'BI', 8)
                pdf.cell(50, 4, "SPED:", ln=1)
                
                pdf.set_font("Helvetica", '', 8)
                for item in diff['sped'][:3]: # Limit to 3 items
                    pdf.set_xy(x_pos, pdf.get_y())
                    safe_item = f"- {item}".encode('latin-1', 'replace').decode('latin-1')
                    pdf.multi_cell(55, 4, safe_item)
                
                y_pos = pdf.get_y() + 2
                
            if diff.get('ml') and y_pos < 250:
                pdf.set_xy(x_pos, y_pos)
                pdf.set_font("Helvetica", 'BI', 8)
                pdf.cell(50, 4, "ML Support:", ln=1)
                
                pdf.set_font("Helvetica", '', 8)
                for item in diff['ml'][:3]: # Limit to 3 items
                    pdf.set_xy(x_pos, pdf.get_y())
                    safe_item = f"- {item}".encode('latin-1', 'replace').decode('latin-1')
                    pdf.multi_cell(55, 4, safe_item)
                
                y_pos = pdf.get_y() + 10
        
        # --- Main Content (Right Side) ---
        # Draw Border Line
        pdf.set_draw_color(200, 200, 200)
        pdf.line(sidebar_width, 38, sidebar_width, 279)
        
        y_pos = 45
        x_pos = sidebar_width + 10 # 74mm
        content_width = 130
        
        for section in data.get('sections', []):
            if y_pos > 260: break # Stop if page full
            
            pdf.set_xy(x_pos, y_pos)
            pdf.set_font("Helvetica", 'B', 11)
            pdf.set_text_color(13, 148, 136) # Teal accent
            phase = section.get('phase', 'Phase')
            time = section.get('time', '')
            pdf.cell(content_width, 6, f"{phase} ({time})", ln=1)
            y_pos += 6
            
            pdf.set_xy(x_pos, y_pos)
            pdf.set_font("Helvetica", '', 10)
            pdf.set_text_color(0, 0, 0)
            activity = section.get('activity', '')
            
            # Truncate if too long (approx 400 chars)
            if len(activity) > 400:
                activity = activity[:397] + "..."
                
            safe_activity = activity.encode('latin-1', 'replace').decode('latin-1')
            pdf.multi_cell(content_width, 5, safe_activity)
            y_pos = pdf.get_y() + 6
            
            y_pos = pdf.get_y() + 6
            
        return pdf.output(dest='S').encode('latin-1'), response.text
        
    except Exception as e:
        st.error(f"Error generating lesson plan: {e}")
        return None, ""

def generate_slide_deck(topic, grade, strategy="None / Standard", source_text=""):
    """Generates a 7-slide PowerPoint presentation using Gemini and python-pptx."""
    # Safety Default
    target_standard = locals().get('target_standard') or globals().get('target_standard') or "General Standard"

    client = get_gemini_client()
    if not client: return None

    # Strategy Context
    strategy_instruction = ""
    if strategy and strategy != "None / Standard":
        strategy_instruction = f"""
        CRITICAL: Include a specific slide titled "{strategy} Activity Instructions".
        On this slide, provide student-facing directions that align exactly with the {strategy} method.
        (e.g., if "Station Rotation", list what happens at each station; if "Fishbowl", list the rules for the inner/outer circle).
        In the Speaker Notes for this slide, provide teacher-facing tips on how to facilitate the activity (e.g., "Set a timer for 10 minutes").
        """
    else:
        strategy_instruction = """
        Include a slide titled "Practice Activity" with clear student instructions for a standard class activity.
        """

    # Prompt Logic: Chain vs Scratch
    if source_text:
        prompt = f"""
        ### 1. ROLE
        Act as an Educational Content Creator specializing in visual presentation design.
        
        ### 2. TASK
        Convert the following Lesson Plan into a 7-slide PowerPoint presentation for Grade {grade}.
        
        ### 3. SOURCE MATERIAL (LESSON PLAN)
        \"\"\"{source_text}\"\"\"
        
        ### 4. CRITICAL: NO TEXT OVERLAP RULE
        The Lesson Plan above contains DETAILED procedural descriptions and activity explanations.
        Your slides must be COMPLEMENTARY, not duplicative:
        
        - Lesson Plan = WHAT students DO (detailed activities, procedures, timing)
        - Slides = VISUAL ANCHORS (key terms, diagrams to reference, keywords ONLY)
        - Speaker Notes = VERBAL DELIVERY (what the teacher says, questions to ask, transitions)
        
        DO NOT copy or paraphrase text from the lesson plan into the slides.
        Instead, extract only the KEY VOCABULARY and CONCEPT NAMES as bullet points.
        
        ### 5. REQUIREMENTS
        - Create 7 slides based on the lesson plan content.
        - {strategy_instruction}
        - Constraint: Use MAXIMUM 5 bullet points per slide.
        - Constraint: Use MAXIMUM 8 words per bullet point. Be extremely concise.
        - Constraint: Move ALL explanations and details into the speaker_notes. The slide text must be keywords only.
        - Constraint: Bullet points should be KEYWORDS/PHRASES, not sentences from the lesson plan.
        
        ### 6. OUTPUT FORMAT
        Return a JSON object with:
        - 'slides': list of objects, where each slide has:
            - 'title': string (Clear and Action-Oriented)
            - 'bullet_points': list of strings (Max 5 points, Max 8 words each. Keywords only.)
            - 'speaker_notes': string (Detailed, scripted notes. e.g., "Ask the class: Have you ever seen...?")
            - 'image_ai_prompt': string (detailed Nano Banana prompt)
        """
    else:
        prompt = f"""
        Create a 7-slide presentation outline for Grade {grade} on "{topic}".
        {strategy_instruction}
        
        Constraint: Use MAXIMUM 5 bullet points per slide.
        Constraint: Use MAXIMUM 8 words per bullet point. Be extremely concise.
        Constraint: Move ALL explanations and details into the speaker_notes. The slide text must be keywords only.
        
        Return a JSON object with:
        - 'slides': list of objects, where each slide has:
            - 'title': string (Clear and Action-Oriented)
            - 'bullet_points': list of strings (Max 5 points, Max 8 words each. Keywords only.)
            - 'speaker_notes': string (Detailed, scripted notes. e.g., "Ask the class: Have you ever seen...?")
            - 'image_ai_prompt': string (detailed Nano Banana prompt)
        """
    
    try:
        response = client.models.generate_content(
            model='gemini-2.0-flash',
            contents=prompt,
            config=types.GenerateContentConfig(response_mime_type='application/json')
        )
        data = clean_json(response.text)
        
        # UNIVERSAL HANDLER: Support List or Dict
        if isinstance(data, list):
            slides_data = data
        elif isinstance(data, dict):
            # Try every possible key the AI might use
            slides_data = data.get('slides') or data.get('presentation') or data.get('content') or list(data.values())[0]
        else:
            slides_data = []
            
        # Safety Check: Ensure it's actually a list
        if not isinstance(slides_data, list):
            slides_data = []
        
        prs = Presentation()
        
        # Helper to add a slide
        def add_slide_content(prs, title_text, bullets_list, notes_text, ai_prompt, is_first=False):
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            
            # Body Content (Standard Placeholder)
            if len(slide.placeholders) > 1:
                content = slide.placeholders[1]
                content.width = Inches(4.5)
                content.height = Inches(5.5)
                content.top = Inches(1.5)
                
                tf = content.text_frame
                tf.word_wrap = True
                tf.clear() # Clear existing
                
                for b in bullets_list:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0

            # Speaker Notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = notes_text

            # Slide 1: Pro Tip
            if is_first:
                left = Inches(0.5)
                top = Inches(0.2) # Very top
                width = Inches(9.0)
                height = Inches(0.5)
                
                tip_box = slide.shapes.add_textbox(left, top, width, height)
                tf = tip_box.text_frame
                p = tf.add_paragraph()
                p.text = "💡 PRO TIP: To style this presentation instantly, click the Design tab and select Designer (or a Theme) to match your classroom style."
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(100, 100, 100) # Grey

            # Footer: Nano Banana Prompt
            left = Inches(0.5)
            top = Inches(7.0)
            width = Inches(9.0)
            height = Inches(0.5)
            
            disc_box = slide.shapes.add_textbox(left, top, width, height)
            tf = disc_box.text_frame
            p = tf.add_paragraph()
            p.text = f"🍌 Nano Banana Image Prompt: {ai_prompt}"
            p.font.italic = True
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(150, 150, 150) # Light Grey

        for i, slide_info in enumerate(slides_data):
            title = slide_info.get('title', 'Untitled Slide')
            bullets = slide_info.get('bullet_points', [])
            notes = slide_info.get('speaker_notes', '')
            prompt = slide_info.get('image_ai_prompt', f"Image of {topic}")
            
            # Split if too many bullets
            if len(bullets) > 6:
                # Part 1
                add_slide_content(prs, f"{title} (Part 1)", bullets[:6], notes, prompt, is_first=(i==0))
                # Part 2
                add_slide_content(prs, f"{title} (Part 2)", bullets[6:], notes, prompt, is_first=False)
            else:
                add_slide_content(prs, title, bullets, notes, prompt, is_first=(i==0))
            
        # Save to buffer
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        return pptx_buffer.getvalue()
        
    except Exception as e:
        st.error(f"Error generating slides: {e}")
        return None

def extract_text_from_file(uploaded_file):
    """Extracts text from PDF, DOCX, or TXT files."""
    try:
        text = ""
        if uploaded_file.type == "application/pdf":
            reader = pypdf.PdfReader(uploaded_file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = docx.Document(uploaded_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif uploaded_file.type == "text/plain":
            text = uploaded_file.getvalue().decode("utf-8")
        
        # Truncate if too long
        if len(text) > 10000:
            text = text[:10000] + "\n\n[Text truncated for prompt limit...]"
            
        return text
    except Exception as e:
        st.error(f"Error reading file: {e}")
        return ""



def construct_quiz_prompt(topic, subtopic, count, due_date, due_time, points_per_question, question_types, grade_level, is_sped, is_gifted, is_ml, language, source_text="", context_topics=None):
    # Build Context Strings
    sped_context = "Include specific accommodations for Special Education (SPED) students." if is_sped else ""
    gifted_context = "Include extension questions and advanced critical thinking challenges for Gifted/Advanced learners." if is_gifted else ""
    ml_context = f"Include language supports for Multilingual Learners (primary language: {language})." if is_ml else ""

    # Contextual Awareness for Unit Mode
    context_instruction = ""
    if context_topics:
        topics_str = ", ".join(context_topics)
        context_instruction = f"CRITICAL: Create a distinct Quiz assessing the following topics covered recently: {topics_str}. Do not re-test older topics."

    # Build Task Constraints
    task_constraints = ""
    if is_sped:
        task_constraints += "Modify reading level to be accessible. Chunk text into smaller sections.\n"
    if is_gifted:
        task_constraints += "Include bonus/challenge questions that require higher-order thinking.\n"
    if is_ml:
        task_constraints += f"Provide key vocabulary definitions translated into {language}.\n"
    # Source Material Section (Protected)
    source_material_section = ""
    if source_text:
        source_material_section = f"""
5. SOURCE MATERIAL
Use the following text as the primary source for content generation: \"\"\" {source_text} \"\"\"
"""

    prompt = f"""
1. ROLE
Act as an expert Curriculum Designer for Grade {grade_level}.

2. CONTEXT
I am teaching a unit on "{topic}" (Subtopic: {subtopic}). Student Profile: Mixed ability. {sped_context} {gifted_context} {ml_context}

3. TASK
Create a Quiz that aligns perfectly with the standard above. {context_instruction} {task_constraints} Generate {count} questions.

XML CONFIGURATION: Set the point value for EVERY question to {points_per_question}. QUESTION TYPES: Generate a mix of ONLY the following types: {question_types}. METADATA: Include the Due Date ({due_date} {due_time}) in the Quiz Description text.

4. FORMAT
Return the output as a JSON object matching the following schema: {{ "questions": [ {{ "question_text": "string", "options": ["string", "string", "string", "string"], "correct_answer_index": int }} ] }} {source_material_section}
"""
    return prompt

def construct_assignment_prompt(topic, subtopic, tool, due_date, due_time, points, grade_level, is_sped, is_gifted, is_ml, language, subject, strategy, source_text=""):
    tools_str = tool if tool and tool != "None" else "None"

    # Build Context Strings
    sped_context = "Include specific accommodations for Special Education (SPED) students." if is_sped else ""
    gifted_context = "Include extension activities and advanced challenges for Gifted/Advanced learners." if is_gifted else ""
    ml_context = f"Include language supports for Multilingual Learners (primary language: {language})." if is_ml else ""
    # Subject Specific Logic
    subject_context = ""
    iframe_height = "450"
    if "Business" in subject or "Economics" in subject:
        subject_context = "Focus on market dynamics, finance, and management scenarios."
        if "Desmos" in tool or "EconGraphs" in tool:
            subject_context += ' Create scenarios that require students to shift the curves. Ask them to predict the new Equilibrium Price.'
            iframe_height = "700"
    elif "Humanities" in subject or "Arts" in subject:
        subject_context = "Focus on history, literature, visual arts, and cultural analysis."
    elif "Technology" in subject or "CS" in subject:
        subject_context = "Focus on coding, digital literacy, and systems thinking."
    # Pedagogy
    pedagogy_section = ""
    if strategy and strategy != "None / Standard":
        pedagogy_section = f"""
5. PEDAGOGY & STRATEGY
Method: {strategy}
"""

    # Source Material
    source_material_section = ""
    if source_text:
        source_material_section = f"""
6. SOURCE MATERIAL
Use the following text as the primary source for content generation: \"\"\" {source_text} \"\"\"
"""

    prompt = f"""
1. ROLE
Act as an expert Curriculum Designer for Grade {grade_level}.

2. CONTEXT
I am teaching a unit on "{topic}" (Subtopic: {subtopic}). Subject: {subject} Student Profile: Mixed ability. {sped_context} {gifted_context} {ml_context}

3. TASK
Create an Assignment that aligns perfectly with the standard above. {subject_context} Tools to embed: {tools_str}

In the HTML output, create a highly visible "Metadata Box" at the top using a styled div that displays: Due Date: {due_date} at {due_time} | Points: {points}.

4. FORMAT
Return the output as raw HTML code ready for Canvas LMS.

Structure: Title, Introduction, Content, Rubric/Answer Key.

Styling: Use inline CSS for a clean, modern look.

Embeds: Include this tool:

{pedagogy_section} {source_material_section}
"""
    return prompt

def construct_unit_prompt(topic, num_assignments, num_quizzes, grade_level, is_sped, is_gifted, is_ml, language, subject, strategy, source_text=""):
    # Build Context Strings
    sped_context = "Include specific accommodations for Special Education (SPED) students." if is_sped else ""
    gifted_context = "Include extension activities and advanced challenges for Gifted/Advanced learners." if is_gifted else ""
    ml_context = f"Include language supports for Multilingual Learners (primary language: {language})." if is_ml else ""
    
    # Source Material
    source_material_section = ""
    if source_text:
        source_material_section = f"""
### 5. SOURCE MATERIAL
Use the following text as the primary source for content generation:
\"\"\" {source_text} \"\"\"
"""
    # The Prompt (Note the f""" wrapper!)
    prompt = f"""
1. ROLE
Act as an expert Curriculum Designer for Grade {grade_level}.

2. CONTEXT
I am planning a comprehensive unit on "{topic}". Subject: {subject} Student Profile: Mixed ability. {sped_context} {gifted_context} {ml_context}

3. TASK
Create a logical unit sequence mixing {num_assignments} Assignments and {num_quizzes} Quizzes. Instructional Strategy: {strategy} Instruction: Place quizzes after relevant assignments to assess learning. Instruction: Create a mixed sequence (e.g., A, A, Q, A, A, Q). Do NOT group all assignments first.

4. OUTPUT FORMAT
Return the output as a JSON list of objects. Each object must have:

"type": "Assignment" or "Quiz"

"title": string (Creative title)

"focus_topic": string (Specific subtopic covered)

Example: [ {{ "type": "Assignment", "title": "Intro to Cells", "focus_topic": "Cell Theory" }}, {{ "type": "Quiz", "title": "Cell Theory Check", "focus_topic": "Cell Theory" }} ]

{source_material_section} """
    return prompt



# --- UI ---

# Professional Header
st.title("Canvas Content Creator")
st.caption("Generate Canvas-ready Assignments, Quizzes, and Full Units with AI.")

st.divider()

# Sidebar
with st.sidebar:
    st.header("Settings")
    

    
    # Unified Content Type Selection
    content_type = st.selectbox("Content Type", ["Assignment", "Quiz", "Unit"])
    
    # File Uploader
    uploaded_file = st.file_uploader("Attach Source Material (PDF, DOCX, TXT)", type=['pdf', 'docx', 'txt'])
    source_text = ""
    if uploaded_file:
        with st.spinner("Extracting text..."):
            source_text = extract_text_from_file(uploaded_file)
            st.success("File processed!")
            
    subject = st.selectbox("Subject Focus", ["Science", "Math", "English", "Business & Economics", "Humanities & Arts", "Technology & CS", "General"])
    topic = st.text_input("Topic", "Photosynthesis")
    
    # Subtopic (only shown for non-Unit content types, but defined here for layout)
    subtopic = ""
    if content_type != 'Unit':
        subtopic = st.text_input("Subtopic", "Light-dependent reactions")
    
    # Standard Input (Needed for Auto-Detection)
    standard = st.text_area("Standard", "NGSS HS-LS1-5", height=100)
    
    # Grade Level
    grade_level = st.selectbox("Grade Level", ['K', '1', '2', '3', '4', '5', '6', '7', '8', '9', '10', '11', '12', 'Higher Ed / Collegiate'], index=10)
    
    # Instructional Strategy
    instructional_strategy = st.selectbox(
        "Instructional Strategy",
        ["None / Standard", "Blended Learning (Station Rotation)", "Project-Based Learning (PBL)", "Flipped Classroom", "Inquiry-Based Learning", "Socratic Seminar / Fishbowl", "Gamification", "Direct Instruction"]
    )
    
    # Media Expansion Packs
    media_packs = st.multiselect(
        "Media Expansion Packs",
        ['Nano Banana (Image Generation)', 'Veo (Video Generation)']
    )
    
    # Conditional Inputs based on Content Type
    num_assignments = 5
    num_quizzes = 2
    
    if content_type == 'Unit':
        st.info("Unit Mode: Generates a comprehensive plan.")
        num_assignments = st.number_input("Number of Assignments", 1, 10, 5)
        num_quizzes = st.number_input("Number of Quizzes", 1, 5, 2)

    st.divider()
    st.subheader("Differentiation")
    is_sped = st.toggle("SPED Accommodations")
    is_gifted = st.toggle("Gifted Extensions")
    is_ml = st.toggle("Multilingual Support")
    
    language = "Spanish"
    if is_ml:
        language = st.selectbox("Target Language", ["Spanish", "French", "Portuguese", "Arabic", "Chinese", "Vietnamese", "Tagalog"])
    
    st.divider()
    st.subheader("Logistics")
    due_date = st.date_input("Due Date", datetime.date.today() + datetime.timedelta(days=7))
    due_time = st.time_input("Due Time", datetime.time(23, 59))
    
    # Logic for Points/Quiz config
    show_assignment_settings = (content_type == 'Assignment')
    show_quiz_settings = (content_type == 'Quiz')
    
    points = 100
    if show_assignment_settings or content_type == 'Unit':
        points = st.number_input("Total Points (Assignments)", value=100)
        
    points_per_question = 1
    question_types = ['Multiple Choice']
    if show_quiz_settings or content_type == 'Unit':
        points_per_question = st.number_input("Points per Question", value=1)
        question_types = st.multiselect("Question Types", options=['Multiple Choice', 'True/False', 'Short Answer', 'Essay', 'Matching', 'Multiple Select'], default=['Multiple Choice'])

    st.header("Tools")
    
    # Auto-Detection Logic
    if 'selected_tool_name' not in st.session_state:
        st.session_state.selected_tool_name = "None"

    if st.button("Auto-Select Best Tool"):
        with st.spinner("Finding the best tool..."):
            recommended = recommend_tool(topic, standard)
            # We check if the recommended tool is valid in general, 
            # but we also need to handle if it's not in the CURRENT subject list.
            if recommended in STEM_TOOLS:
                st.session_state.selected_tool_name = recommended
                st.toast(f"Found match: {recommended}", icon="✅")
            else:
                st.toast("No perfect match found.", icon="⚠️")

    tool_options = ["None"] + list(STEM_TOOLS.keys())
    
    # Filter tools based on Subject
    if subject == "Business & Economics":
        tool_options = ["None", "EconGraphs: Competitive Market", "Desmos: Supply & Demand Shifters", "Marginal Revolution: Elasticity Practice", "Omni Margin Calculator"]
    elif subject == "Humanities & Arts":
        tool_options = ["None", "AutoDraw", "Sketchpad", "Color Wheel", "Google Arts & Culture"]
    elif subject == "Technology & CS":
        tool_options = ["None", "Python Online Compiler", "Scratch"]
    elif subject == "Science":
        # Filter for Science tools + General
        tool_options = ["None"] + [k for k in STEM_TOOLS.keys() if "PhET" in k or "Science" in k or "National Geographic" in k or k in ["YouTube: Crash Course", "YouTube: Khan Academy", "Wikipedia", "Google Slides", "Canva"]]
    elif subject == "Math":
        # Filter for Math tools + General
        tool_options = ["None"] + [k for k in STEM_TOOLS.keys() if "Desmos" in k or "GeoGebra" in k or k in ["YouTube: Khan Academy", "Wikipedia", "Google Slides", "Canva"]]
    
    # Validation: Ensure the currently selected tool is actually in the filtered list.
    # If not, reset to "None".
    if st.session_state.selected_tool_name not in tool_options:
        st.session_state.selected_tool_name = "None"
    
    selected_tool = st.selectbox(
        "Embed Interactive Tool",
        options=tool_options,
        key="selected_tool_name"
    )

# Main Area

if content_type == "Assignment":
    st.header("Assignment Builder")
    
    st.markdown("**Optional Downloads:**")
    col1, col2 = st.columns(2)
    with col1:
        cb1_col, label1_col = st.columns([0.15, 0.85])
        with cb1_col:
            include_lesson_plan = st.checkbox("LP", key="lesson_plan_cb", label_visibility="collapsed")
        with label1_col:
            st.markdown('<p style="color: #3C453C; margin-top: 5px;">Lesson Plan PDF</p>', unsafe_allow_html=True)
    with col2:
        cb2_col, label2_col = st.columns([0.15, 0.85])
        with cb2_col:
            include_slides = st.checkbox("SL", key="slides_cb", label_visibility="collapsed")
        with label2_col:
            st.markdown('<p style="color: #3C453C; margin-top: 5px;">Slides PPTX</p>', unsafe_allow_html=True)
    
    if st.button("Draft My Mega-Prompt"):
        # 1. Generate Main Prompt
        st.session_state['generated_prompt'] = construct_assignment_prompt(topic, subtopic, selected_tool, due_date, due_time, points, grade_level, is_sped, is_gifted, is_ml, language, subject, instructional_strategy, source_text)
        
        # 2. Generate PDF (if checked)
        lp_text = ""
        if include_lesson_plan:
            with st.spinner("Generating Lesson Plan PDF..."):
                st.session_state['lesson_plan_pdf'], lp_text = generate_lesson_plan_pdf(topic, standard, grade_level, instructional_strategy)
        else:
            st.session_state['lesson_plan_pdf'] = None

        # 3. Generate Slides (if checked)
        if include_slides:
            with st.spinner("Generating PowerPoint Slides..."):
                # Chain: Pass the Lesson Plan text if available
                st.session_state['slide_deck_pptx'] = generate_slide_deck(topic, grade_level, instructional_strategy, source_text=lp_text)
        else:
            st.session_state['slide_deck_pptx'] = None
        
        st.session_state['is_generated'] = True

elif content_type == "Quiz":
    st.header("Quiz Builder")
    question_count = st.number_input("Number of Questions", 1, 50, 5)
    
    if st.button("Draft My Mega-Prompt"):
        # 1. Construct Prompt
        prompt_content = construct_quiz_prompt(topic, subtopic, question_count, due_date, due_time, points_per_question, question_types, grade_level, is_sped, is_gifted, is_ml, language, source_text)
        st.session_state['generated_prompt'] = prompt_content
        
        # 2. Generate JSON & Zip (Background)
        # We use the batched function now
        quiz_data = generate_quiz_data_batched(topic, subtopic, question_count, due_date, due_time, points_per_question, question_types, grade_level, is_sped, is_gifted, is_ml, language, source_text)
        
        if quiz_data:
            zip_bytes = generate_qti_zip(quiz_data, title=f"{topic} Quiz")
            st.session_state['quiz_zip'] = zip_bytes
        else:
            st.session_state['quiz_zip'] = None

        st.session_state['is_generated'] = True
        # Reset other artifacts
        st.session_state['lesson_plan_pdf'] = None
        st.session_state['slide_deck_pptx'] = None

elif content_type == "Unit":
    st.header("Unit Planner")
    
    if st.button("Draft My Mega-Prompt"):
        # 1. Generate Sequence
        with st.spinner("Planning Unit Sequence..."):
            prompt = construct_unit_prompt(topic, num_assignments, num_quizzes, grade_level, is_sped, is_gifted, is_ml, language, subject, instructional_strategy, source_text)
            sequence_data = generate_unit_sequence_json(prompt)
            
            # Store the prompt for display
            st.session_state['generated_prompt'] = prompt
            
        if sequence_data:
            # Display the generated sequence
            st.subheader("📋 Generated Unit Sequence")
            for i, item in enumerate(sequence_data):
                item_type = item.get('type', 'Unknown')
                title = item.get('title', f'Item {i+1}')
                focus = item.get('focus_topic', topic)
                icon = "📝" if item_type == "Assignment" else "❓"
                st.markdown(f"{icon} **{i+1}. {item_type}:** {title} *(Focus: {focus})*")
            
            st.markdown("---")
            st.info("📦 Each Assignment will include: HTML file, Lesson Plan PDF, and PowerPoint Slides")
            
            # 2. Generate Package (including Lesson Plans and Slides for each Assignment)
            with st.spinner("Generating Unit Resources (This may take a moment)..."):
                unit_zip = generate_unit_package(
                    sequence_data, topic, grade_level, is_sped, is_gifted, is_ml, language, 
                    subject, instructional_strategy, source_text, due_date, due_time, 
                    points, points_per_question, question_types, standard=standard
                )
                st.session_state['unit_zip'] = unit_zip
        else:
            st.error("Failed to plan unit sequence.")
            st.session_state['unit_zip'] = None

        st.session_state['is_generated'] = True
        # Reset other artifacts
        st.session_state['lesson_plan_pdf'] = None
        st.session_state['slide_deck_pptx'] = None
        st.session_state['quiz_zip'] = None

# Display Results (Persistent)
if st.session_state['is_generated']:
    st.code(st.session_state['generated_prompt'], language='markdown')
    
    # Media Prompts (Explicit Logic)
    if media_packs:
        st.markdown("---")
        
        if "Nano Banana (Image Generation)" in media_packs:
            st.subheader("🍌 Nano Banana Prompt")
            image_prompt = f"Create a 4K educational poster for {topic}. Style: Photorealistic/Diagram. Key elements: [Insert Standard Details]."
            st.code(image_prompt, language='text')
            
        if "Veo (Video Generation)" in media_packs:
            st.subheader("🎥 Veo Video Prompt")
            video_prompt = f"Cinematic 60s video clip. Subject: {topic}. Action: [Describe motion]. Style: Documentary."
            st.code(video_prompt, language='text')
    
    # Download Buttons
    if st.session_state.get('lesson_plan_pdf'):
        st.download_button(
            label="📄 Download Lesson Plan PDF",
            data=st.session_state['lesson_plan_pdf'],
            file_name=f"Lesson_Plan_{topic.replace(' ', '_')}.pdf",
            mime="application/pdf"
        )
    
    if st.session_state.get('slide_deck_pptx'):
        st.download_button(
            label="📊 Download PowerPoint Slides",
            data=st.session_state['slide_deck_pptx'],
            file_name=f"Slides_{topic.replace(' ', '_')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    if st.session_state.get('quiz_zip'):
        st.download_button(
            label="📦 Download Ready-to-Import Quiz (.zip)",
            data=st.session_state['quiz_zip'],
            file_name=f"{topic.replace(' ', '_')}_Quiz.zip",
            mime="application/zip",
            type="primary"
        )
        
    if st.session_state.get('unit_zip'):
        st.download_button(
            label="📦 Download Full Unit Package (.zip)",
            data=st.session_state['unit_zip'],
            file_name=f"Unit_{topic.replace(' ', '_')}.zip",
            mime="application/zip",
            type="primary"
        )
